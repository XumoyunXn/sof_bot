import telebot
from telebot import types
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt, Inches
from googleapiclient.discovery import build
import os, json, re, requests

# 🔑 Kalitlaringiz
TELEGRAM_TOKEN = "7222740719:AAFUSUmIUYGQRPOMB5dwneqZqY-0WwjlH50"
GEMINI_API_KEY = "AIzaSyDm2iGQ-WhiCnika8DgHLsNZ4N4yqQd3q0"
GOOGLE_CSE_ID = "20c10ecfea9a64099"   # Google Custom Search Engine ID
GOOGLE_API_KEY = "AIzaSyBWJn5ya_dVCaCy8wel61CnB-VgUcs-9VQ"   # Google API Key

bot = telebot.TeleBot(TELEGRAM_TOKEN)

# AI sozlash
genai.configure(api_key=GEMINI_API_KEY)
model = genai.GenerativeModel("gemini-1.5-flash")

# User state
user_state = {}

# 🔹 Fayl nomini xavfsiz qilish funksiyasi
def safe_filename(text):
    return re.sub(r'[^a-zA-Z0-9_\-]', '_', text)[:50]

# 🔹 Google qidiruv funksiyasi (faqat http/https linklar qaytaradi)
def search_images(query, num=3):
    service = build("customsearch", "v1", developerKey=GOOGLE_API_KEY)
    res = service.cse().list(q=query, cx=GOOGLE_CSE_ID, searchType="image", num=num).execute()
    links = []
    for item in res.get("items", []):
        link = item.get("link", "")
        if link.startswith("http://") or link.startswith("https://"):
            links.append(link)
    return links


def create_ppt(slides, title, template_file, filename="taqdimot.pptx"):
    prs = Presentation(template_file)

    # Bosh slayd
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "AI tomonidan yaratilgan taqdimot"

    # Slayd o‘lchamlari
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Qolgan slaydlar
    for sl in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # 🔹 Matn joylashuvi (chap tomonda yarim slayd)
        left = Inches(0.5)
        top = Inches(1)
        width = slide_width / 2 - Inches(1)
        height = slide_height - Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        # Sarlavha va matn
        p_title = tf.add_paragraph()
        p_title.text = sl['title']
        p_title.font.bold = True
        p_title.font.size = Pt(24)

        p_content = tf.add_paragraph()
        p_content.text = sl['content']
        p_content.font.size = Pt(18)

        # 🔹 Rasm joylashuvi (o‘ng tomonda yarim slayd)
        try:
            query = f"{sl['title']} {sl['content']}"
            imgs = search_images(query, num=3)

            for img_url in imgs:
                try:
                    if not (img_url.startswith("http://") or img_url.startswith("https://")):
                        continue

                    # Kengaytma aniqlash
                    ext = img_url.split(".")[-1].split("?")[0].lower()
                    if ext not in ["jpg", "jpeg", "png", "webp"]:
                        ext = "jpg"

                    img_data = requests.get(img_url, timeout=10).content
                    img_path = f"temp_{safe_filename(sl['title'])}.{ext}"

                    with open(img_path, "wb") as f:
                        f.write(img_data)

                    # Slaydga rasm joylash
                    img_left = slide_width / 2 + Inches(0.2)
                    img_top = Inches(1)
                    img_width = slide_width / 2 - Inches(1)
                    img_height = slide_height - Inches(2)

                    slide.shapes.add_picture(img_path, img_left, img_top, img_width, img_height)
                    os.remove(img_path)
                    break  # ✅ Agar rasm qo‘shildi — boshqa linklarni tekshirmaydi
                except Exception as inner_e:
                    print(f"⚠️ Rasm ishlamadi, keyingisini sinayapman... {inner_e}")
                    continue
        except Exception as e:
            print("❌ Umumiy rasm qo‘shishda xatolik:", e)

    prs.save(filename)
    return filename


@bot.message_handler(commands=['start'])
def start(message):
    markup = types.ReplyKeyboardMarkup(resize_keyboard=True)
    btn1 = types.KeyboardButton("📑 Taqdimot")
    markup.add(btn1)
    bot.send_message(message.chat.id, "Assalomu alaykum!\nTaqdimot yaratishni xohlaysizmi?", reply_markup=markup)


@bot.message_handler(func=lambda m: m.text == "📑 Taqdimot")
def ask_design(message):
    chat_id = message.chat.id
    user_state[chat_id] = {"step": "choose_design"}

    # Dizayn tanlash
    markup = types.InlineKeyboardMarkup()
    markup.add(types.InlineKeyboardButton("🟢 Slayd Dizayn 1", callback_data="design_white"))
    markup.add(types.InlineKeyboardButton("🔵 Slayd Dizayn 2", callback_data="design_blue"))
    markup.add(types.InlineKeyboardButton("🟣 Slayd Dizayn 3", callback_data="design_green"))

    bot.send_message(chat_id, "Slayd dizaynini tanlang:", reply_markup=markup)


@bot.callback_query_handler(func=lambda call: call.data.startswith("design_"))
def handle_design(call):
    design = call.data.replace("design_", "")
    chat_id = call.message.chat.id

    templates = {
        "white": "templates/test_bot.pptx",
        "blue": "templates/test_bot_2.pptx",
        "green": "templates/test_bot_3.pptx"
    }

    user_state[chat_id]["design"] = templates[design]
    user_state[chat_id]["step"] = "write_title"

    bot.send_message(chat_id, "Endi taqdimot sarlovhasini yozing ✍️")


@bot.message_handler(func=lambda message: True)
def handle_title(message):
    chat_id = message.chat.id

    if chat_id in user_state and user_state[chat_id].get("step") == "write_title":
        title = message.text
        user_state[chat_id]["title"] = title
        user_state[chat_id]["step"] = "generate"

        bot.send_message(chat_id, "⏳ AI taqdimotni yozmoqda, kuting...")

        # AI prompt
        prompt = f"""
        Siz faqat JSON chiqaring.
        Format:
        [
          {{"title": "Slayd sarlavhasi", "content": "Slayd mazmuni"}}
        ]

        Mavzu: {title}
        6 ta slayd yozing.
        """

        response = model.generate_content(prompt)
        raw = response.text

        # JSON tozalash
        match = re.search(r"\[.*\]", raw, re.DOTALL)
        if match:
            json_text = match.group(0)
            slides = json.loads(json_text)
        else:
            bot.send_message(chat_id, "⚠️ AI JSON qaytarmadi:\n\n" + raw)
            return

        # PowerPoint yaratish
        filename = f"taqdimot_{chat_id}.pptx"
        create_ppt(slides, title, user_state[chat_id]["design"], filename)

        with open(filename, "rb") as f:
            bot.send_document(chat_id, f)

        os.remove(filename)


print("🤖 Bot ishlayapti...")
bot.infinity_polling()
